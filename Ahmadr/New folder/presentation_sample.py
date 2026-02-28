from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_slides():
    prs = Presentation()

    # --- HELPER FUNCTION TO ADD TITLE & CONTENT ---
    def add_slide(title_text, content_lines):
        slide_layout = prs.slide_layouts[1] # Bullet Layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        
        # Content
        tf = slide.placeholders[1].text_frame
        tf.text = content_lines[0] # First line
        
        for line in content_lines[1:]:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0

    # ==========================================
    # SLIDE 1: METHODOLOGY (EXTRACTION)
    # ==========================================
    add_slide(
        "Methodology: Phase 1 (Automated Extraction)",
        [
            "Step 1: Data Normalization",
            "   - Created Digital Terrain Model (DTM) to flatten slope terrain.",
            "   - Normalized height (Z) = Raw Z - Local Ground Z.",
            "",
            "Step 2: Canopy Modeling (CHM)",
            "   - Converted 3D Point Cloud to 2D Image.",
            "   - Grid Resolution: 0.1 meters (High Precision).",
            "",
            "Step 3: Segmentation (The Core Logic)",
            "   - Applied 'Watershed Algorithm' to separate overlapping crowns.",
            "   - Used Local Maxima Filter (0.7m radius) to find tree tops."
        ]
    )

    # ==========================================
    # SLIDE 2: METHODOLOGY (AI PREDICTION)
    # ==========================================
    add_slide(
        "Methodology: Phase 2 (AI Prediction)",
        [
            "Feature Engineering:",
            "   - Input: Tree Height (m) & Crown Diameter (m).",
            "   - Extracted automatically from ULS data.",
            "",
            "Machine Learning Model:",
            "   - Algorithm: Random Forest Regressor.",
            "   - Why? Robust against noise and handles non-linear relationships.",
            "",
            "Training Data:",
            "   - Integrated Ground Truth (TLS) + Airborne Data (ALS).",
            "   - Used temporal matching to track growth over 6 months."
        ]
    )

    # ==========================================
    # SLIDE 3: VALIDATION RESULTS (BR01 Plot)
    # ==========================================
    add_slide(
        "Results & Validation (Success Metrics)",
        [
            "Test Area: Plot BR01 (Dense Forest Patch)",
            "",
            "🎯 Detection Accuracy: 96.08%",
            "   - Ground Truth Trees: 51",
            "   - Automatically Detected: 49",
            "",
            "📏 Precision Metrics:",
            "   - Spatial Error: ~15 cm (After Shift Correction)",
            "   - Height Error (MAE): < 1.0 meter (After Calibration)",
            "",
            "⚡ Efficiency:",
            "   - Manual Survey: ~3-4 Hours",
            "   - AI Processing: ~35 Seconds"
        ]
    )

    # ==========================================
    # SLIDE 4: CHALLENGES & SOLUTIONS
    # ==========================================
    add_slide(
        "Technical Challenges & Solutions",
        [
            "🔴 Challenge 1: Sloped Terrain",
            "   - Issue: Standard height calculation gave errors on slopes.",
            "   - ✅ Solution: Implemented Dynamic DTM for local ground referencing.",
            "",
            "🔴 Challenge 2: Ghost Trees (Over-segmentation)",
            "   - Issue: Large branches detected as separate trees.",
            "   - ✅ Solution: Optimized Smoothing Sigma (1.0) & Min-Distance (0.7m).",
            "",
            "🔴 Challenge 3: GPS/LiDAR Misalignment",
            "   - Issue: 1-2 meter shift between datasets.",
            "   - ✅ Solution: Developed Auto-Calibration Algorithm to align coordinates."
        ]
    )

    # Save File
    prs.save('FYP_Addon_Slides.pptx')
    print("✅ Slides Generated! Open 'FYP_Addon_Slides.pptx'")

if __name__ == "__main__":
    create_slides()